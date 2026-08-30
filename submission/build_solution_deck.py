"""Build a slide version of the Detailed Solution Document.

Same content, same numbers and the same palette as the PDF, laid out for a screen rather
than a page. Every figure here is the one the PDF carries, which is the one the artifacts
under docs/results/ produce.

    python submission/build_solution_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
REPO = HERE.parent
OUT = HERE / "ControlPlane_Detailed_Solution.pptx"

INK = RGBColor(0x17, 0x12, 0x1F)
BODY = RGBColor(0x2E, 0x28, 0x38)
MUTED = RGBColor(0x6B, 0x64, 0x78)
ACCENT = RGBColor(0x6D, 0x28, 0xD9)
BRIGHT = RGBColor(0xA1, 0x00, 0xFF)
DEEP = RGBColor(0x3B, 0x07, 0x64)
RULE = RGBColor(0xE3, 0xDE, 0xEC)
PANEL = RGBColor(0xF8, 0xF6, 0xFC)
PANEL2 = RGBColor(0xF1, 0xEC, 0xFA)
AMBER = RGBColor(0x9A, 0x5B, 0x00)
AMBERBG = RGBColor(0xFF, 0xF7, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def text_box(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, spacing=1.0):
    """runs is a list of (text, size, bold, colour, space_after_pt)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    for index, (text, size, bold, colour, after) in enumerate(runs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(after)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = FONT
    return box


def rect(slide, x, y, w, h, fill, line=None, line_w=0.9):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def slide_base(prs, number, title, kicker=None, footer="Detailed Solution Document"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, WHITE)
    if kicker:
        text_box(slide, Inches(0.75), Inches(0.42), Inches(11), Inches(0.25),
                 [(kicker.upper(), 10, True, ACCENT, 0)])
    text_box(slide, Inches(0.75), Inches(0.68), Inches(11.8), Inches(0.55),
             [(title, 27, True, INK, 0)])
    rect(slide, Inches(0.75), Inches(1.32), Inches(11.83), Pt(2.2), ACCENT)
    text_box(slide, Inches(11.9), Inches(6.94), Inches(0.7), Inches(0.25),
             [(str(number), 10, True, MUTED, 0)], align=PP_ALIGN.RIGHT)
    text_box(slide, Inches(0.75), Inches(6.94), Inches(6), Inches(0.25),
             [(f"ControlPlane  .  {footer}", 9, False, MUTED, 0)])
    return slide


def table(slide, x, y, w, headers, rows, widths, *, size=11, header_size=10, row_h=0.34):
    """A hand drawn table. python-pptx tables fight the theme, so this stays in shapes."""
    total = sum(widths)
    cols = [Emu(int(w * part / total)) for part in widths]
    rect(slide, x, y, w, Inches(row_h), PANEL2)
    cursor = x
    for index, head in enumerate(headers):
        text_box(slide, cursor + Inches(0.09), y + Inches(0.075), cols[index] - Inches(0.14),
                 Inches(0.24), [(head, header_size, True, INK, 0)])
        cursor += cols[index]
    rect(slide, x, y + Inches(row_h), w, Pt(1.4), ACCENT)
    top = y + Inches(row_h) + Pt(1.4)
    for r, row in enumerate(rows):
        if r % 2 == 1:
            rect(slide, x, top, w, Inches(row_h), RGBColor(0xFB, 0xFA, 0xFE))
        cursor = x
        for c, cell in enumerate(row):
            bold = cell.startswith("*")
            label = cell[1:] if bold else cell
            text_box(slide, cursor + Inches(0.09), top + Inches(0.06), cols[c] - Inches(0.14),
                     Inches(0.26), [(label, size, bold, INK if bold else BODY, 0)])
            cursor += cols[c]
        rect(slide, x, top + Inches(row_h), w, Pt(0.6), RULE)
        top += Inches(row_h) + Pt(0.6)
    return top


def stats(slide, y, items, *, x=Inches(0.75), w=Inches(11.83)):
    gap = Inches(0.22)
    each = Emu(int((w - gap * (len(items) - 1)) / len(items)))
    cursor = x
    for value, label in items:
        rect(slide, cursor, y, each, Pt(3), ACCENT)
        text_box(slide, cursor, y + Inches(0.14), each, Inches(0.45),
                 [(value, 24, True, INK, 0)])
        text_box(slide, cursor, y + Inches(0.62), each, Inches(0.6),
                 [(label, 11, False, MUTED, 0)], spacing=1.15)
        cursor += each + gap


def note(slide, y, title, body, *, tone="take", x=Inches(0.75), w=Inches(11.83), h=Inches(0.86)):
    fill = PANEL if tone == "take" else AMBERBG
    edge = ACCENT if tone == "take" else AMBER
    rect(slide, x, y, w, h, fill)
    rect(slide, x, y, Pt(3.4), h, edge)
    text_box(slide, x + Inches(0.22), y + Inches(0.13), w - Inches(0.44), h - Inches(0.24),
             [(title, 12, True, DEEP if tone == "take" else AMBER, 3),
              (body, 12, False, BODY, 0)], spacing=1.16)


def bullets(slide, x, y, w, items, *, size=13, gap=8):
    runs = []
    for head, rest in items:
        runs.append((head, size, True, INK, 1))
        runs.append((rest, size, False, BODY, gap))
    text_box(slide, x, y, w, Inches(4), runs, spacing=1.2)


def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------- 1 cover ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, WHITE)
    rect(s, 0, 0, Inches(0.34), H, DEEP)
    text_box(s, Inches(1.1), Inches(1.15), Inches(10), Inches(0.3),
             [("ACCENTURE INNOVATION CHALLENGE 2026  .  ROUND 2  .  TRACK 1", 12, True, ACCENT, 0)])
    text_box(s, Inches(1.1), Inches(1.72), Inches(10), Inches(1.1),
             [("ControlPlane", 60, True, INK, 0)])
    text_box(s, Inches(1.1), Inches(2.82), Inches(10), Inches(0.5),
             [("Detailed Solution Document", 25, False, ACCENT, 0)])
    rect(s, Inches(1.1), Inches(3.46), Inches(2.1), Pt(3.4), BRIGHT)
    text_box(s, Inches(1.1), Inches(3.78), Inches(9.4), Inches(1.1),
             [("An assurance gateway that decides how much checking each AI answer is worth, "
               "enforces a safety floor the budget cannot override, and writes every decision "
               "into a ledger that can be recomputed from its own arithmetic.", 15, False, BODY, 0)],
             spacing=1.3)
    stats(s, Inches(5.25), [
        ("92.8%", "of the protection of checking every answer, for 10.0% of the compute"),
        ("1,500", "held out interactions behind every figure"),
        ("0.008", "to 0.043 calibration error by route"),
        ("1,500 of 1,500", "decisions in one verified chain"),
    ], x=Inches(1.1), w=Inches(11.4))
    text_box(s, Inches(1.1), Inches(6.86), Inches(11), Inches(0.3),
             [("Team BurningShadows   .   github.com/Advait-Pathak9222/BurningShadows   .   "
               "Reproduce with ./run_submission.sh   .   Offline, no API key, no network, no GPU",
               10, False, MUTED, 0)])

    # ---------- 2 contents ----------
    s = slide_base(prs, 2, "Contents", "What follows")
    left = [("01  Executive summary", "What this is and what it demonstrates"),
            ("02  The problem", "Why fixed verification wastes money in both directions"),
            ("03  Solution approach", "Verification as a spending decision"),
            ("04  Architecture", "Four planes and a learning loop"),
            ("05  End to end request flow", "One held out request, stage by stage"),
            ("06  Core decision logic", "Tiers, floor, allocation, review, effects"),
            ("07  Audit and reproducibility", "The ledger and the regeneration path")]
    right = [("08  Implementation", "What is built, and where it lives"),
             ("09  Evaluation and evidence", "Measured results on held out traffic"),
             ("10  External benchmarks", "Five public corpora, including one we failed"),
             ("11  Running the system", "Commands, console and live prototype"),
             ("12  Deployment approach", "Where each tier belongs in production"),
             ("13  Validation boundaries", "What is proven and what is not"),
             ("14  Conclusion", "What we would do next")]
    for column, items in ((Inches(0.75), left), (Inches(6.95), right)):
        top = Inches(1.75)
        for head, rest in items:
            text_box(s, column, top, Inches(5.6), Inches(0.28), [(head, 14, True, INK, 0)])
            text_box(s, column, top + Inches(0.27), Inches(5.6), Inches(0.24),
                     [(rest, 11, False, MUTED, 0)])
            rect(s, column, top + Inches(0.58), Inches(5.6), Pt(0.6), RULE)
            top += Inches(0.72)

    # ---------- 3 problem ----------
    s = slide_base(prs, 3, "The problem", "01  .  02")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.5),
             [("Verification of AI output is treated as a fixed rule when it is really a "
               "spending decision under a budget.", 16, False, BODY, 0)], spacing=1.25)
    cards = [
        ("Checking everything", "The dearest tier costs 160 times the cheapest and adds 225 times "
         "the delay. On two million interactions a year that is Rs 64 lakh of compute before a "
         "single reviewer is paid."),
        ("Sampling blindly", "A one in ten sample spends the same effort on a question about "
         "opening hours as on an instruction to transfer funds, because it cannot tell them apart."),
        ("No defensible record", "When a regulator asks why one answer was released, most stacks "
         "produce a log line. What is needed is the arithmetic, and it is usually gone."),
    ]
    x = Inches(0.75)
    for title, body in cards:
        rect(s, x, Inches(2.35), Inches(3.78), Inches(1.72), WHITE, RULE)
        rect(s, x, Inches(2.35), Inches(3.78), Pt(3), ACCENT)
        text_box(s, x + Inches(0.24), Inches(2.55), Inches(3.3), Inches(0.3),
                 [(title, 15, True, DEEP, 0)])
        text_box(s, x + Inches(0.24), Inches(2.92), Inches(3.3), Inches(1.05),
                 [(body, 11.5, False, BODY, 0)], spacing=1.2)
        x += Inches(4.03)
    note(s, Inches(4.35), "The bill is mostly people. ",
         "One completed review costs Rs 120 against Rs 3.20 for the dearest automated check, so a "
         "review is 37.5 times an automated check. Across the budget grid, 81 to 98 paise in every "
         "assurance rupee is human review time rather than compute.", h=Inches(0.95))
    table(s, Inches(0.75), Inches(5.52), Inches(11.83),
          ["Compute budget", "Automated checking", "Human review", "Attention share", "Cases raised"],
          [["10%", "Rs 480.98", "Rs 19,920", "*97.6%", "246"],
           ["40%", "Rs 1,920.18", "Rs 19,920", "91.2%", "344"],
           ["100%", "Rs 4,800.00", "Rs 19,920", "80.6%", "569"]],
          [2.2, 2.4, 2.2, 2.4, 2.0], size=11.5)

    # ---------- 4 approach ----------
    s = slide_base(prs, 4, "Solution approach", "03")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.5),
             [("Treat every check as a purchase. Estimate what it would prevent, compare that with "
               "what it costs, and buy it only when the first number is larger. Then put a floor "
               "underneath so safety is never a purely economic question.", 15, False, BODY, 0)],
             spacing=1.25)
    rect(s, Inches(0.75), Inches(2.5), Inches(11.83), Inches(0.92), PANEL, RULE)
    text_box(s, Inches(1.1), Inches(2.66), Inches(11), Inches(0.6),
             [("expected loss  =  calibrated risk  x  consequence", 15, False, DEEP, 4),
              ("check when   expected loss x catch rate   >   (1 + shadow price) x check cost",
               15, True, DEEP, 0)], spacing=1.15)
    bullets(s, Inches(0.75), Inches(3.72), Inches(5.7), [
        ("Harm is a vector, not a label. ", "An answer can leak data and be wrong at once, so five "
         "axes are kept separate through the whole pipeline."),
        ("Evidence regime decides what is knowable. ", "An unverifiable answer carrying real risk "
         "is withheld rather than guessed at."),
        ("The floor is not negotiable. ", "A release threshold selected on held out data. The budget "
         "can slow everything else. It cannot switch this off."),
    ], size=12.5)
    bullets(s, Inches(6.85), Inches(3.72), Inches(5.7), [
        ("Attention is the scarce resource. ", "Reviewer minutes cost 37.5 times an automated check "
         "and cannot be bought at short notice, so the queue is scheduled by value against deadline."),
        ("Words and actions are different risks. ", "Text can stream while checking runs because it "
         "can be retracted. A payment cannot, so effects wait behind their own gate."),
    ], size=12.5)

    # ---------- 5 architecture ----------
    s = slide_base(prs, 5, "Architecture", "04")
    text_box(s, Inches(0.75), Inches(1.58), Inches(11.83), Inches(0.36),
             [("The gateway sees only the request, the response, any supplied context and any "
               "proposed tool calls. It needs no model weights, hidden states or log probabilities.",
               13, False, BODY, 0)])
    s.shapes.add_picture(str(REPO / "docs" / "images" / "architecture.png"),
                         Inches(0.75), Inches(2.12), width=Inches(11.83))
    note(s, Inches(6.05), "Why this shape. ",
         "The allocator never needs to know what a tier is. It needs three numbers per tier, which "
         "are cost, catch rate and latency. Swapping a detector is an adapter implementation rather "
         "than a rewrite.", h=Inches(0.78))

    # ---------- 6 request flow ----------
    s = slide_base(prs, 6, "End to end request flow", "05")
    text_box(s, Inches(0.75), Inches(1.58), Inches(11.83), Inches(0.36),
             [("Held out row cp-02477. A finance route request carrying a transfer, where the model "
               "echoes a secret the source said must never be repeated. Every figure is recomputed "
               "from the stored record.", 13, False, BODY, 0)])
    table(s, Inches(0.75), Inches(2.12), Inches(11.83),
          ["#", "Stage", "What happened to this request"],
          [["1", "*A request arrives", "Route finops-agent, jurisdiction india, with a transfer_funds call attached"],
           ["2", "*Policy is resolved", "Pack india-dpdp-2025.11. Exfiltration on this route is priced at Rs 55,000"],
           ["3", "*The model answers", "Two cheap tiers read it for 20 paise. The answer repeats a secret from the context"],
           ["4", "*Scores become probabilities", "Prompt attack 1.000, hallucination 0.250, private data 0.150. Regime grounded"],
           ["5", "*The safety floor is consulted", "Route threshold 0.25 against highest harm 1.000. The check is mandatory"],
           ["6", "*Every tier is priced", "Loss on the table Rs 69,357. Tier 2 prevents Rs 63,123 for Rs 3.20 and is chosen"],
           ["7", "*A verdict is reached", "With the judge merged in, harm is 0.910, above the 0.88 block boundary. Blocked"],
           ["8", "*The action is gated apart", "transfer_funds, classified financial, denied. The money never moves"],
           ["9", "*A receipt is written", "Tier sums, threshold, shadow price and policy hash, chained to the record before it"]],
          [0.5, 3.1, 8.4], size=11.5, row_h=0.36)
    note(s, Inches(6.05), "Stages five and six are the whole product. ",
         "The floor decides whether a check is optional. The allocator decides which check is worth "
         "buying. Everything else keeps those two decisions honest and recomputable.", h=Inches(0.78))

    # ---------- 7 tiers and floor ----------
    s = slide_base(prs, 7, "Core decision logic  .  tiers and the release floor", "06.1  .  06.2")
    text_box(s, Inches(0.75), Inches(1.6), Inches(5.7), Inches(0.3),
             [("Three tiers behind one adapter contract", 14, True, DEEP, 0)])
    table(s, Inches(0.75), Inches(2.0), Inches(5.7),
          ["Tier", "Cost", "Latency", "Multiple", "Catch, leak"],
          [["*0  Rules", "Rs 0.02", "4 ms", "1x", "0.82"],
           ["*1  Small models", "Rs 0.18", "70 ms", "9x", "0.90"],
           ["*2  Judge", "Rs 3.20", "900 ms", "160x", "0.94"]],
          [2.0, 1.2, 1.2, 1.1, 1.4], size=11.5)
    text_box(s, Inches(0.75), Inches(3.66), Inches(5.7), Inches(0.9),
             [("Catch rates are measured, not configured. On the held out run the judge caught "
               "0.930 against 0.880 configured over 365 observations, so a better detector reports "
               "its own higher catch rate and the allocator buys more of it.", 12, False, BODY, 0)],
             spacing=1.2)
    text_box(s, Inches(6.85), Inches(1.6), Inches(5.73), Inches(0.3),
             [("The floor, validated on held out data", 14, True, DEEP, 0)])
    table(s, Inches(6.85), Inches(2.0), Inches(5.73),
          ["Route", "Threshold", "Released", "Escaped", "Rate", "Bound"],
          [["finops-agent", "0.25", "339", "19", "*0.0560", "0.0756"],
           ["internal-kb", "1.00", "476", "34", "*0.0714", "0.0890"],
           ["support-assistant", "0.10", "397", "21", "*0.0529", "0.0703"]],
          [2.3, 1.4, 1.2, 1.1, 1.2, 1.2], size=11)
    text_box(s, Inches(6.85), Inches(3.66), Inches(5.73), Inches(0.9),
             [("Every route holds well inside a bound of 0.15 at 90% confidence, and holds non "
               "vacuously, meaning rows really were released unchecked and the bound still bound.",
               12, False, BODY, 0)], spacing=1.2)
    note(s, Inches(4.75), "The floor is only informative when the bound exceeds the harm base rate. ",
         "An operator whose traffic is one third harmful at a bound of 0.15 is not getting a "
         "guarantee, they are getting full coverage, and should either raise the bound or budget "
         "for checking everything. That rule is computable before deployment.", h=Inches(0.92))

    # ---------- 8 allocation and review ----------
    s = slide_base(prs, 8, "Core decision logic  .  budget and human review", "06.3  .  06.4  .  06.5")
    text_box(s, Inches(0.75), Inches(1.6), Inches(11.83), Inches(0.34),
             [("A shadow price alone cannot bound spend. Ours did not, until the governor reserved "
               "the floor cost.", 13.5, False, BODY, 0)])
    stats(s, Inches(2.02), [
        ("1.00x to 1.03x", "spend against budget across the grid, once governed, against up to 3.75x before"),
        ("Rs 51.84", "what the floor alone obliges, 1.08% of full coverage"),
        ("5.625%", "below this budget fraction, blanket cheap coverage stops being affordable"),
    ])
    text_box(s, Inches(0.75), Inches(3.3), Inches(11.83), Inches(0.3),
             [("Order alone decides which cases a saturated review desk loses", 14, True, DEEP, 0)])
    table(s, Inches(0.75), Inches(3.72), Inches(11.83),
          ["Serving rule", "Expected loss served", "Deadlines missed", "High value cases shed"],
          [["*deadline_density, shipped", "*Rs 3,430,681", "65", "*2"],
           ["density, ablation", "Rs 3,798,647", "48", "2"],
           ["random, baseline", "Rs 2,816,476", "47", "9"],
           ["fifo, baseline", "Rs 2,162,224", "139", "15"]],
          [3.6, 3.0, 2.6, 2.6], size=12)
    note(s, Inches(5.72), "1.59x more expected loss from the same 166 reviews. ",
         "High value cases dropped fall from 15 to 2. The density ablation leads on both axes and is "
         "reported as the stronger rule rather than hidden. Ordering is still the smaller lever, "
         "because keeping up needs 3.0 reviewers against the two staffed.", h=Inches(0.92))

    # ---------- 9 audit ----------
    s = slide_base(prs, 9, "Audit and reproducibility", "07")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.36),
             [("An assurance system that cannot show its working is an opinion. Two mechanisms make "
               "this one checkable, one at runtime and one at build time.", 15, False, BODY, 0)])
    stats(s, Inches(2.2), [
        ("1,500 of 1,500", "decisions recorded in one chain"),
        ("224 of 224", "proposed effects logged"),
        ("205", "reviews in the same chain"),
        ("Valid", "verification over 1,705 records"),
    ])
    bullets(s, Inches(0.75), Inches(3.85), Inches(5.7), [
        ("The ledger. ", "Each record carries the hash of the one before it. It holds the calibrated "
         "vector, the raw scores, the tier arithmetic, the threshold, the shadow price and the "
         "policy hash, which is enough to recompute the decision rather than recall it."),
    ], size=12.5)
    bullets(s, Inches(6.85), Inches(3.85), Inches(5.7), [
        ("The regeneration path. ", "./run_submission.sh builds a clean environment, runs the gate, "
         "regenerates every artifact and fails if one differs from the committed copy. If a number "
         "stops reproducing, the script names the file."),
    ], size=12.5)
    note(s, Inches(5.5), "Endpoints were fixed before the work started. ",
         "Results are reported against those criteria whether or not they were met. The record "
         "includes a queue defect written down before it was corrected, and a case where fitting a "
         "detector on the rows that certified its bound made it claim 0.1407 while held out data "
         "showed 0.2800. The discipline is the guarantee.", h=Inches(1.05))

    # ---------- 10 evaluation ----------
    s = slide_base(prs, 10, "Evaluation and evidence", "09")
    text_box(s, Inches(0.75), Inches(1.6), Inches(11.83), Inches(0.32),
             [("1,500 held out interactions across three routes and two jurisdictions, with span "
               "level ground truth. Four policies across six budgets.", 13.5, False, BODY, 0)])
    table(s, Inches(0.75), Inches(2.08), Inches(11.83),
          ["Policy", "Compute spend", "Loss averted", "Share of benefit", "Share of cost", "Escaped harm"],
          [["Check nothing", "Rs 0.00", "Rs 3,184,800", "58.2%", "0.0%", "0.0962"],
           ["*Allocator at 10% budget", "*Rs 480.98", "*Rs 5,078,000", "*92.8%", "*10.0%", "0.0363"],
           ["Allocator at 25% budget", "Rs 1,098.58", "Rs 5,407,000", "98.9%", "22.9%", "0.0211"],
           ["Allocator at 80% budget", "Rs 3,839.64", "Rs 5,477,300", "100.1%", "80.0%", "0.0117"],
           ["Check everything", "Rs 4,800.00", "Rs 5,469,400", "100.0%", "100.0%", "0.0120"]],
          [3.0, 2.1, 2.2, 1.9, 1.7, 1.6], size=11.5)
    note(s, Inches(4.5), "At 80% it beats checking everything. ",
         "It averts more loss for 20% less compute, because checking everything spends its budget on "
         "rows where the judge had nothing to find.", h=Inches(0.72))
    stats(s, Inches(5.5), [
        ("0.008 to 0.043", "expected calibration error across the three routes"),
        ("10.9%", "of tier choices change across a 0.25x to 4x price swing"),
        ("0%", "of verdicts flip across that same swing"),
    ])

    # ---------- 11 honest results ----------
    s = slide_base(prs, 11, "Three results reported because they went against us", "09  .  10")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.34),
             [("Read the null column before the score. On an imbalanced corpus a policy that flags "
               "every row can look respectable while detecting nothing.", 14, False, BODY, 0)])
    table(s, Inches(0.75), Inches(2.12), Inches(11.83),
          ["Benchmark", "ControlPlane", "Flag everything null", "Margin", "Published comparison", "Verdict"],
          [["ToxicChat, AUPRC", "*0.597", "0.071", "+0.526", "Llama Guard 0.664, OpenAI Mod 0.588", "In band"],
           ["RAGTruth, F1", "*0.601", "0.518", "+0.083", "LettuceDetect 0.792, GPT-4 0.634", "In band"],
           ["BeaverTails at 7% harm, F1", "0.165", "0.139", "+0.026", "None on our label mapping", "Margin only"],
           ["*Aegis, AUPRC", "*0.811", "0.661", "+0.151", "Band 0.860 to 0.941", "*Below band, failed"],
           ["OR-Bench, AUC", "0.784", "0.500", "+0.284", "Operating point compared separately", "Partial"]],
          [3.0, 1.7, 2.2, 1.3, 3.7, 1.9], size=11)
    bullets(s, Inches(0.75), Inches(4.55), Inches(11.83), [
        ("Allocation is not universally better than a well tuned fixed rate policy. ",
         "At matched spend it wins 5 of 7 budgets on our corpus and only 1 of 7 on Aegis. It helps "
         "when the harm mix varies and blanket cheap coverage is unaffordable, and both conditions "
         "are computable before deployment."),
        ("The learning loop currently refuses on all three routes. ",
         "A single pass over 1,500 rows yields 73 usable labelled pairs against roughly 100 per "
         "route needed, so make relearn declines to release and says why. That is the gate working."),
    ], size=12.5)

    # ---------- 12 boundaries ----------
    s = slide_base(prs, 12, "Current validation boundaries", "13")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.34),
             [("Being straight about this is more useful than claiming readiness. This is what a "
               "production deployment would have to close.", 14, False, BODY, 0)])
    table(s, Inches(0.75), Inches(2.12), Inches(11.83),
          ["Boundary", "Severity", "What it means and what closes it"],
          [["*Calibration fitted on synthetic traffic", "High",
            "Maps and thresholds do not transfer. Roughly 500 labelled rows per route, about four weeks of one reviewer"],
           ["*Drift breaks the bound silently", "High",
            "The guarantee is valid only for the distribution it was calibrated on. Drift detection is not built"],
           ["*Tier 2 is a deterministic stand in", "Medium",
            "The adapter exists and two real adapters are written against it. A production judge must be built and re measured"],
           ["*Single process state", "Medium", "The ledger and budget controller assume one process. Multiple workers need shared state"],
           ["*Effect gate has no durable lease", "Medium", "It decides and records. It does not hold a lock across a crash"],
           ["*Consequence prices need finance sign off", "Low", "Getting this wrong costs money rather than safety, because price does not enter the release rule"]],
          [3.9, 1.3, 8.6], size=11, row_h=0.4)
    note(s, Inches(5.55), "Most reviewer labels cannot be learned from. ",
         "The queue serves by value, so harmful rows are likelier to be reviewed within the raised "
         "population. That is selection inside a stratum and no stratum level weight undoes it. The "
         "sample runs 37.2% harmful against 18.4% in traffic, and weighting moved it only to 38.2%. "
         "That is why one reviewer slot in five is filled at random.", tone="warn", h=Inches(1.1))

    # ---------- 13 conclusion ----------
    s = slide_base(prs, 13, "Conclusion", "14")
    text_box(s, Inches(0.75), Inches(1.75), Inches(11.83), Inches(1.0),
             [("ControlPlane demonstrates that verification of AI output can be run as a priced "
               "decision rather than a fixed rule. Doing so is worth roughly 90% of the protection "
               "of checking everything for a tenth of the compute, while holding a safety bound the "
               "budget cannot switch off and producing a record that can be recomputed rather than "
               "trusted.", 16, False, BODY, 0)], spacing=1.3)
    text_box(s, Inches(0.75), Inches(3.05), Inches(11.83), Inches(0.8),
             [("It also demonstrates the limits of that claim honestly. Allocation helps under "
               "conditions computable in advance and does not help outside them. The detectors are "
               "stubs, which is what makes the allocation result attributable and is the first thing "
               "a deployment would replace. The reviewer queue, not the compute budget, is where the "
               "money and the risk actually sit.", 14, False, BODY, 0)], spacing=1.25)
    note(s, Inches(4.35), "The next step is the one the evidence points at. ",
         "Take one route in one enterprise, run the gateway in shadow mode beside the existing "
         "stack, put a real judge behind the Tier 2 adapter, and refit calibration on that route's "
         "own labelled traffic until the learning loop stops refusing to release. Everything needed "
         "to do that is in this repository.", h=Inches(1.1))
    text_box(s, Inches(0.75), Inches(5.85), Inches(11.83), Inches(0.4),
             [("github.com/Advait-Pathak9222/BurningShadows      .      make demo      .      "
               "make prototype      .      ./run_submission.sh", 13, True, ACCENT, 0)])

    prs.save(OUT)
    print(f"{OUT.name}  {len(prs.slides.__iter__.__self__._sldIdLst)} slides  "
          f"{OUT.stat().st_size / 1024 / 1024:.2f} MB")


if __name__ == "__main__":
    build()
