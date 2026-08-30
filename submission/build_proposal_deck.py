"""Build a slide version of the Business Proposal.

Same content, same numbers and the same palette as the proposal PDF, laid out for a room
rather than a page. Layout helpers are shared with the solution deck so the two decks and
the two PDFs stay one visual family.

    python submission/build_proposal_deck.py
"""

from __future__ import annotations

from pathlib import Path

from build_solution_deck import (
    ACCENT,
    BODY,
    BRIGHT,
    DEEP,
    INK,
    MUTED,
    RULE,
    WHITE,
    H,
    W,
    bullets,
    note,
    rect,
    slide_base,
    stats,
    table,
    text_box,
)
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "ControlPlane_Business_Proposal.pptx"


def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------- 1 cover ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, WHITE)
    rect(s, 0, 0, Inches(0.34), H, DEEP)
    text_box(s, Inches(1.1), Inches(1.15), Inches(10), Inches(0.3),
             [("ACCENTURE INNOVATION CHALLENGE 2026  .  ROUND 2  .  TRACK 1", 12, True, ACCENT, 0)])
    text_box(s, Inches(1.1), Inches(1.72), Inches(10), Inches(1.1),
             [("ControlPlane", 60, True, INK, 0)])
    text_box(s, Inches(1.1), Inches(2.82), Inches(10), Inches(0.5),
             [("Business Proposal", 25, False, ACCENT, 0)])
    rect(s, Inches(1.1), Inches(3.46), Inches(2.1), Pt(3.4), BRIGHT)
    text_box(s, Inches(1.1), Inches(3.78), Inches(9.6), Inches(1.1),
             [("Enterprises are putting AI in front of customers, staff and money. The cost of "
               "checking what it says is becoming a line item nobody planned for, and most of that "
               "line item is people. This is a proposal to spend it deliberately.",
               15, False, BODY, 0)], spacing=1.3)
    stats(s, Inches(5.25), [
        ("81 to 98%", "of the assurance bill is human review time, not compute"),
        ("37.5x", "what one review costs against the dearest automated check"),
        ("1.59x", "more expected loss handled from the same reviewer hours"),
        ("8 weeks", "to a decision, on one route, with no enforcement risk"),
    ], x=Inches(1.1), w=Inches(11.4))
    text_box(s, Inches(1.1), Inches(6.86), Inches(11), Inches(0.3),
             [("Team BurningShadows   .   Prepared for enterprise AI platform and risk owners   .   "
               "The ask is one route, one shadow pilot, eight weeks", 10, False, MUTED, 0)])

    # ---------- 2 executive summary ----------
    s = slide_base(prs, 2, "Executive summary", "The case in one slide", footer="Business Proposal")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.8),
             [("Every enterprise deploying AI assistants has quietly acquired a new operating cost. "
               "Somebody has to check what the model says before it reaches a customer or moves "
               "money. That cost is growing faster than the AI spend it protects, and almost nobody "
               "is managing it.", 16, False, BODY, 0)], spacing=1.28)
    box = s.shapes.add_textbox(Inches(0.75), Inches(2.66), Inches(11.83), Inches(0.85))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    para = frame.paragraphs[0]
    para.line_spacing = 1.28
    for text, bold, colour in (
        ("ControlPlane prices the difference. ", True, INK),
        ("It estimates what it would cost the business if a given answer were wrong, and buys "
         "checking only where the loss prevented is worth more than the check. A safety floor sits "
         "underneath that the budget cannot switch off.", False, BODY),
    ):
        run = para.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Segoe UI"
    stats(s, Inches(3.85), [
        ("92.8%", "of the protection of checking every answer, for 10% of the compute"),
        ("4.29x", "more loss left on the table by doing nothing at all"),
        ("2 of 15", "high value cases dropped by our review order, against 15 by first in first out"),
        ("0%", "of verdicts change when every price assumption is swung four fold"),
    ])
    note(s, Inches(5.55), "What we are asking for. ",
         "One customer facing route, in shadow mode, for eight weeks. No enforcement, no latency "
         "risk, no change to the existing model. At the end you will have your own numbers for what "
         "assurance costs you today and what it would cost run this way.", h=Inches(0.95))

    # ---------- 3 problem ----------
    s = slide_base(prs, 3, "Problem framing", "Why this has no owner today", footer="Business Proposal")
    text_box(s, Inches(0.75), Inches(1.62), Inches(11.83), Inches(0.4),
             [("The problem is not that AI makes mistakes. It is that checking for mistakes has no "
               "budget, no owner and no way to prove it worked.", 16, False, BODY, 0)])
    cards = [
        ("Overspending", "Checking everything", "Adds 900 milliseconds to every answer and, at two "
         "million interactions a year, costs Rs 64 lakh in compute before a single reviewer is paid."),
        ("Underspending", "Sampling blindly", "A fixed sample cannot tell a question about opening "
         "hours from an instruction to move money, so it protects the cheap and misses the expensive."),
        ("Unprovable", "No defensible record", "When a regulator asks why one answer was released, a "
         "log line is not an answer. The arithmetic behind the decision is what is needed."),
    ]
    x = Inches(0.75)
    for tag, title, body in cards:
        rect(s, x, Inches(2.3), Inches(3.78), Inches(1.85), WHITE, RULE)
        rect(s, x, Inches(2.3), Inches(3.78), Pt(3), ACCENT)
        text_box(s, x + Inches(0.24), Inches(2.5), Inches(3.3), Inches(0.22),
                 [(tag.upper(), 9.5, True, ACCENT, 0)])
        text_box(s, x + Inches(0.24), Inches(2.78), Inches(3.3), Inches(0.3),
                 [(title, 15, True, DEEP, 0)])
        text_box(s, x + Inches(0.24), Inches(3.15), Inches(3.3), Inches(1.0),
                 [(body, 11.5, False, BODY, 0)], spacing=1.2)
        x += Inches(4.03)
    text_box(s, Inches(0.75), Inches(4.42), Inches(11.83), Inches(0.3),
             [("The cost nobody planned for", 15, True, DEEP, 0)])
    table(s, Inches(0.75), Inches(4.82), Inches(11.83),
          ["What it costs to check one answer", "Price", "Who pays it"],
          [["Rules and pattern scan", "Rs 0.02", "Compute budget"],
           ["Small model checks", "Rs 0.18", "Compute budget"],
           ["Language model judge", "Rs 3.20", "Compute budget"],
           ["*One completed human review", "*Rs 120.00", "*Headcount budget"]],
          [5.6, 3.0, 3.3], size=12)

    # ---------- 4 the finding ----------
    s = slide_base(prs, 4, "The finding that shapes this proposal", "Business opportunity", footer="Business Proposal")
    text_box(s, Inches(0.75), Inches(1.7), Inches(11.83), Inches(0.9),
             [("We set out to save compute. The measurement said something more useful.",
               17, False, BODY, 6),
              ("Between 81 and 98 paise in every assurance rupee is human review time, not machine "
               "time.", 21, True, DEEP, 0)], spacing=1.25)
    text_box(s, Inches(0.75), Inches(2.95), Inches(11.83), Inches(0.5),
             [("Every vendor in this category competes on the 2 to 19 paise. Nobody is managing the "
               "rest. That changes what is worth buying, because review capacity is fixed and the "
               "queue is already oversubscribed.", 15, False, BODY, 0)], spacing=1.25)
    text_box(s, Inches(0.75), Inches(3.7), Inches(11.83), Inches(0.3),
             [("Where the value sits", 15, True, DEEP, 0)])
    table(s, Inches(0.75), Inches(4.1), Inches(11.83),
          ["Value", "What it is worth", "Confidence today"],
          [["*Reviewer hours redirected", "The same capacity handles 1.59 times more expected loss, and drops 2 high value cases instead of 15", "Measured on our corpus"],
           ["*Compute avoided", "92.8% of the protection of checking everything, for 10% of the compute", "Measured on our corpus"],
           ["*Incidents avoided", "Doing nothing leaves 4.29 times more residual loss on the same traffic", "Depends on your consequence prices"],
           ["*Audit cost avoided", "Every decision recomputable from its own record, with the policy version attached", "Built and verified"]],
          [3.1, 7.0, 2.8], size=11.5, row_h=0.38)
    note(s, Inches(6.0), "What we are not claiming. ",
         "The rupee value of loss averted is a property of the consequence prices we assumed, not "
         "evidence of a customer saving. Those prices are the first thing a pilot replaces with "
         "yours.", tone="warn", h=Inches(0.78))

    # ---------- 5 solution ----------
    s = slide_base(prs, 5, "Proposed solution", "What it does, in the order it does it", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(11.83),
          ["Step", "What happens"],
          [["*Refuse early", "A cheap scan reads the prompt before the model runs. An obvious attack is refused there and costs neither generation nor checking"],
           ["*Score the answer", "Cheap checks score the answer on five kinds of harm for a combined 20 paise, and those scores become probabilities"],
           ["*Apply the floor", "Each route has a safety threshold above which a check is compulsory. The budget cannot switch this off"],
           ["*Price the choice", "Below the floor, the system buys the expensive check only where the loss it would prevent beats what it costs"],
           ["*Split words from actions", "Text can be shown while checking runs. Actions that move money or change records wait behind their own gate"],
           ["*Prioritise people", "Whatever needs a person is queued by value per reviewer minute against the deadline, not first come first served"],
           ["*Record everything", "Every decision, its inputs and its price go into a chained record where a later edit is detectable"]],
          [2.9, 10.0], size=12, row_h=0.38)
    note(s, Inches(5.3), "What integration actually involves. ",
         "The gateway presents the interface your application already calls, so integration is a "
         "change of base URL. It reads only the request, the answer, any supplied context and any "
         "proposed actions, so it needs no access to model internals and can run inside your own "
         "network where data residency requires it.", h=Inches(1.0))

    # ---------- 6 buyers ----------
    s = slide_base(prs, 6, "Target users and buyers", "Four stakeholders, four questions", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.7), Inches(11.83),
          ["Stakeholder", "Decision they own", "What this answers for them"],
          [["*Economic buyer  .  CFO or FinOps", "The assurance budget and reviewer headcount",
            "Total cost of assurance split into compute and attention, with cost per thousand interactions reported"],
           ["*Risk owner  .  Chief Risk", "The minimum control and audit evidence",
            "A per route safety floor with a stated bound, jurisdiction policy packs, every override in a verifiable chain"],
           ["*Technical owner  .  AI Platform", "Integration, portability and reliability",
            "A drop in gateway on the interface you already use, detector adapters that are one method, load tested"],
           ["*Operational user  .  Head of Review", "Reviewer workload and what gets dropped",
            "A queue ordered by value against deadline, with shed rate and deadline misses reported rather than hidden"]],
          [3.5, 3.4, 6.9], size=11.5, row_h=0.46)
    text_box(s, Inches(0.75), Inches(4.35), Inches(11.83), Inches(0.5),
             [("The budget usually already exists as manual review headcount rather than as a "
               "tooling line. That matters for how this is sold. It is a reallocation of an existing "
               "cost, not a new one.", 14, False, BODY, 0)], spacing=1.25)
    note(s, Inches(5.15), "Where it fits first. ",
         "A regulated enterprise running a customer facing assistant alongside an assistant that can "
         "take an action. Banking, insurance, telecom and large scale consumer support all have the "
         "same shape. High volume, a real consequence table, an existing review desk, and a "
         "regulator who will eventually ask for evidence.", h=Inches(1.0))

    # ---------- 7 use case ----------
    s = slide_base(prs, 7, "Customer use case", "Two requests, one minute apart, same assistant", footer="Business Proposal")
    for index, (tag, title, body, cost) in enumerate([
        ("Request one", "A routine ticket lookup",
         "A customer asks for the status of an open ticket. The cheap checks score the answer well "
         "below the safety threshold. The expensive judge would prevent very little extra loss and "
         "would cost far more than that, so it is not bought. The answer is released.",
         "Cost of assurance 20 paise.   Delay added, none the customer notices."),
        ("Request two", "A payment instruction",
         "A request on the finance route carries a transfer, and the model echoes a secret its "
         "source said must never be repeated. Risk lands above the safety threshold, so a check is "
         "compulsory before any price is considered. The judge is bought, the answer is blocked, "
         "the transfer is denied, and the case is queued for a person.",
         "Cost of assurance Rs 3.20.   Loss on the table Rs 69,357."),
    ]):
        x = Inches(0.75) + index * Inches(6.13)
        rect(s, x, Inches(1.7), Inches(5.7), Inches(3.1), WHITE, RULE)
        rect(s, x, Inches(1.7), Inches(5.7), Pt(3), ACCENT)
        text_box(s, x + Inches(0.28), Inches(1.92), Inches(5.1), Inches(0.22),
                 [(tag.upper(), 9.5, True, ACCENT, 0)])
        text_box(s, x + Inches(0.28), Inches(2.2), Inches(5.1), Inches(0.32),
                 [(title, 17, True, DEEP, 0)])
        text_box(s, x + Inches(0.28), Inches(2.66), Inches(5.1), Inches(1.5),
                 [(body, 12.5, False, BODY, 0)], spacing=1.22)
        text_box(s, x + Inches(0.28), Inches(4.3), Inches(5.1), Inches(0.3),
                 [(cost, 12, True, INK, 0)])
    note(s, Inches(5.05), "The point of the pair. ",
         "Neither request was treated as average. The first was cheap because nothing was at stake. "
         "The second was expensive because a great deal was, and the system could afford it "
         "precisely because the first one did not spend.", h=Inches(0.85))
    text_box(s, Inches(0.75), Inches(6.1), Inches(11.83), Inches(0.5),
             [("What the customer never sees. ", 13, True, INK, 0)])
    text_box(s, Inches(3.0), Inches(6.1), Inches(9.6), Inches(0.5),
             [("In the second case the customer receives a message saying the request has been "
               "passed to a colleague. The business does not lose the money, and the review desk "
               "receives a case tagged with what it is worth.", 13, False, BODY, 0)], spacing=1.2)

    # ---------- 8 business case ----------
    s = slide_base(prs, 8, "Business case", "Every input is exposed so it can be challenged", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(11.83),
          ["Input", "Where it comes from", "Status today"],
          [["*Reviewer rate and handling time", "Invoice and queue study", "Defensible. This is a real number a CFO already owns"],
           ["*Check price", "Provider token prices", "Defensible once a real provider adapter is in place"],
           ["*Catch rate", "Reviewer labels plus a sampled audit", "Measured. The judge caught 0.930 against 0.880 configured over 365 observations"],
           ["*Consequence of each harm", "Finance and risk, as a range", "Assumption, but bounded. Across a 0.25 to 4 times swing, 10.9% of spending decisions move and no verdict does"]],
          [3.4, 3.2, 7.2], size=11.5, row_h=0.42)
    stats(s, Inches(3.72), [
        ("Rs 64 lakh", "a year in compute to check every answer at two million interactions"),
        ("Rs 6.4 lakh", "a year for 92.8% of the same protection"),
        ("81 to 98%", "of the remaining bill is reviewer time, which is the real target"),
    ])
    note(s, Inches(5.3), "Stated plainly, because it is the strongest objection. ",
         "At the tightest budget a well tuned blanket cheap check averts 2.4% more than we do for a "
         "fraction of the compute spend. Our advantage on compute alone is small. Where we beat the "
         "obvious alternative is on the resource that costs the money, which is reviewer time.",
         tone="warn", h=Inches(1.0))

    # ---------- 9 impact and differentiation ----------
    s = slide_base(prs, 9, "Expected impact and differentiation", "What changes, and why it holds", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(11.83),
          ["Where", "Effect measured on held out traffic"],
          [["*Compute spend", "92.8% of the protection of checking everything for 10.0% of the cost. At an 80% budget it beats checking everything"],
           ["*Reviewer productivity", "1.59 times more expected loss from the same hours. High value cases dropped fall from 15 to 2. Deadline misses fall from 139 to 65"],
           ["*Customer experience", "Text is never held behind a verdict, only actions are. Under overload the effect tail holds at 104 ms against 1,508 ms"],
           ["*Audit readiness", "1,500 of 1,500 decisions and 224 of 224 proposed actions logged in one verified chain"]],
          [2.9, 10.0], size=11.5, row_h=0.4)
    text_box(s, Inches(0.75), Inches(3.95), Inches(11.83), Inches(0.3),
             [("We are not competing with detectors. We are the layer that decides which detector to "
               "buy, and when.", 15, True, DEEP, 0)])
    table(s, Inches(0.75), Inches(4.35), Inches(11.83),
          ["Approach", "What it does", "What it leaves unanswered"],
          [["Guardrail and moderation tools", "Classify text as safe or unsafe, usually well", "How much to spend on this answer"],
           ["LLM as judge platforms", "Apply a strong model to check output", "Which requests are worth that price"],
           ["Observability and evaluation", "Report what happened after the fact", "Anything at request time. They do not gate an action"],
           ["*ControlPlane", "*Prices the check, enforces a floor, orders the queue and proves the decision", "*Detection quality, which is why any detector plugs in"]],
          [3.4, 5.4, 4.4], size=11, row_h=0.36)
    note(s, Inches(6.15), "The defensible position. ",
         "Better detectors help us rather than threaten us. A better detector reports a higher catch "
         "rate and the allocator automatically buys more of it.", h=Inches(0.62))

    # ---------- 10 commercial and pilot ----------
    s = slide_base(prs, 10, "Commercial approach and pilot plan", "Eight weeks, one route, no enforcement risk", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(5.7),
          ["Element", "Proposal"],
          [["*Deployment", "Self hosted beside your gateway. Integration is a base URL change"],
           ["*Pricing", "Per thousand interactions assessed, with the budget set by you"],
           ["*Land", "One customer facing route in shadow mode"],
           ["*Expand", "Add the route that can act, then enforce on capped effects"],
           ["*Value narrative", "Reviewer hours redirected, not GPU spend saved"]],
          [1.9, 3.8], size=11, row_h=0.42)
    table(s, Inches(6.85), Inches(1.68), Inches(5.73),
          ["Phase", "What happens", "What you get"],
          [["*Weeks 1 to 2", "Shadow beside the existing stack. No enforcement", "Your own baseline cost, split compute and attention"],
           ["*Weeks 3 to 4", "Fit calibration and the threshold on your traffic", "A release bound stated on your data"],
           ["*Weeks 5 to 6", "Compare fixed rate against allocation offline", "An honest answer on whether allocation helps here"],
           ["*Weeks 7 to 8", "Holds for capped effects, with human approval", "A working control and a signed runbook"]],
          [1.6, 2.4, 2.4], size=10.5, row_h=0.48)
    text_box(s, Inches(0.75), Inches(4.4), Inches(11.83), Inches(0.3),
             [("Promotion to production requires all of these", 14, True, DEEP, 0)])
    bullets(s, Inches(0.75), Inches(4.78), Inches(11.83), [
        ("", "Allocation beats the fixed rate baseline at equal spend across agreed budget bands, "
             "with intervals rather than point estimates."),
        ("", "The route release bound stays under the approved threshold on fresh labels, and added "
             "latency at the 99th percentile stays inside the route service level objective."),
        ("", "Every proposed action has an audit record, decisions stay stable across the approved "
             "consequence ranges, and the incident owner has approved the runbook."),
    ], size=12, gap=5)

    # ---------- 11 metrics and roadmap ----------
    s = slide_base(prs, 11, "Success metrics and roadmap", "How you will know, and what comes next", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(11.83),
          ["Metric", "What it tells you", "Target for the pilot"],
          [["*Cost per thousand interactions", "Total assurance cost, split into compute and attention", "Below your current blended cost"],
           ["*Escaped harm rate", "Share of released answers that carried harm", "Under the approved bound on fresh labels"],
           ["*Expected loss served per reviewer hour", "Whether review capacity goes to the right cases", "Above the first in first out baseline"],
           ["*High value cases shed", "Which cases you lose when capacity runs out", "Materially below today's rate"],
           ["*Audit completeness", "Share of decisions and actions with a verifiable record", "100 percent"]],
          [4.0, 5.2, 4.0], size=11, row_h=0.36)
    table(s, Inches(0.75), Inches(4.35), Inches(11.83),
          ["Phase", "Scope", "Evidence that closes the phase"],
          [["*Now  .  Prototype", "Offline gateway, allocation, floor, queue, effect gating, audit chain", "Reproducible commands and an honest baseline. Complete"],
           ["*Phase 1  .  Shadow pilot", "Real detector adapters, provider adapter, sampled audit slice", "Catch rates and calibration fitted on your route"],
           ["*Phase 2  .  Controlled effects", "Reviewer identity, capped held actions, durable effect leases", "No unaudited actions, approved override workflow"],
           ["*Phase 3  .  Scale", "Shared calibration, signed policy releases, drift detection", "Multi worker load and recovery tests"]],
          [3.0, 5.4, 4.8], size=11, row_h=0.38)

    # ---------- 12 risks and boundaries ----------
    s = slide_base(prs, 12, "Risks and current validation boundaries", "Stated before you ask", footer="Business Proposal")
    table(s, Inches(0.75), Inches(1.68), Inches(11.83),
          ["Risk", "Mitigation and stop condition"],
          [["*Consequence prices steer the wrong traffic", "Measured. 10.9% of spending decisions move across a 0.25 to 4 times band against a 20% stop condition, and no verdict moves"],
           ["*Reviewer capacity is the real constraint", "Confirmed, and it is. Attention is 81 to 98% of total cost, which is why queue ordering is part of the product"],
           ["*Selective labels inflate the catch rate", "A fixed rate audit of released rows plus a random sampling reserve, both implemented and costed"],
           ["*Traffic drift invalidates the bound", "Window monitoring and label refresh. Drift detection is Phase 3 and is not built today"],
           ["*The safety floor exceeds the budget", "Report infeasibility. The floor wins and the overspend is declared. The bound is never relaxed silently"]],
          [4.2, 9.0], size=11.5, row_h=0.38)
    text_box(s, Inches(0.75), Inches(4.35), Inches(11.83), Inches(0.3),
             [("What would have to close before production", 14, True, DEEP, 0)])
    bullets(s, Inches(0.75), Inches(4.72), Inches(11.83), [
        ("", "Calibration is fitted on synthetic traffic and does not transfer. Roughly 500 labelled "
             "rows per route, about four weeks of one reviewer, would refit it."),
        ("", "The expensive judge is a deterministic stand in. The adapter exists and two real "
             "adapters are written against it, but a production judge must be built and re measured."),
        ("", "Drift is not detected today, the effect gate does not hold a lock across a crash, and "
             "the record store assumes a single process."),
    ], size=12, gap=5)

    # ---------- 13 close ----------
    s = slide_base(prs, 13, "Closing recommendation", "The ask", footer="Business Proposal")
    text_box(s, Inches(0.75), Inches(1.85), Inches(11.83), Inches(0.6),
             [("Approve an eight week shadow pilot on one customer facing route, with no "
               "enforcement.", 22, True, DEEP, 0)], spacing=1.2)
    text_box(s, Inches(0.75), Inches(2.75), Inches(11.83), Inches(1.0),
             [("It costs almost nothing and it settles the argument with your own data. In two weeks "
               "you will know what assurance actually costs you today, split between compute and "
               "people, which most organisations cannot currently state. In eight weeks you will "
               "know whether allocating that spend beats what you do now on your traffic, and you "
               "will have a release bound measured on your own labels rather than on ours.",
               16, False, BODY, 0)], spacing=1.28)
    note(s, Inches(4.1), "We will also tell you if the answer is no. ",
         "Allocation helps under conditions computable in advance, and does not help outside them. "
         "On one of the five public corpora we tested it won on only one budget setting out of "
         "seven, and we report that alongside the results where it won. A pilot that concludes the "
         "honest answer is a fixed rate policy on your traffic is a successful pilot, and it will "
         "have cost you eight weeks and no enforcement risk to find out.", h=Inches(1.25))
    rect(s, Inches(0.75), Inches(5.72), Inches(11.83), Inches(0.72), DEEP)
    text_box(s, Inches(1.1), Inches(5.92), Inches(11.2), Inches(0.4),
             [("One route.      One shadow pilot.      Eight weeks.      "
               "At the end you own the numbers.", 18, True, WHITE, 0)], align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"{OUT.name}  {len(prs.slides._sldIdLst)} slides  "
          f"{OUT.stat().st_size / 1024 / 1024:.2f} MB")


if __name__ == "__main__":
    build()
